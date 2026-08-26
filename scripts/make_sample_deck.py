#!/usr/bin/env python3
"""Generate the sample deck served on the upload page (public/sample-deck.pptx).
Re-run this script whenever you want to change the sample."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BLUE = RGBColor(0x1A, 0x73, 0xE8)
DARK = RGBColor(0x20, 0x21, 0x24)
GRAY = RGBColor(0x5F, 0x63, 0x68)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1 — Title
s = prs.slides.add_slide(prs.slide_layouts[6])
tb = s.shapes.add_textbox(Inches(1), Inches(2.4), Inches(11.3), Inches(1.4))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "JPEIGO Slides"
r.font.size = Pt(54); r.font.bold = True; r.font.color.rgb = BLUE
tb2 = s.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.3), Inches(1))
p2 = tb2.text_frame.paragraphs[0]
r2 = p2.add_run(); r2.text = "Sample deck for trying the translator"
r2.font.size = Pt(24); r2.font.color.rgb = GRAY

# Slide 2 — Why translate with us
s = prs.slides.add_slide(prs.slide_layouts[6])
tb = s.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.3), Inches(1))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Why teams choose JPEIGO"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = DARK
points = [
    "Formatting stays exactly as designed",
    "Tables, charts and text boxes preserved",
    "Fast AI translation with review controls",
    "Download ready-to-present files",
]
tb2 = s.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.9), Inches(4))
tf = tb2.text_frame
for i, pt_text in enumerate(points):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(18)
    r = p.add_run(); r.text = "•  " + pt_text
    r.font.size = Pt(26); r.font.color.rgb = DARK

# Slide 3 — Table demo
s = prs.slides.add_slide(prs.slide_layouts[6])
tb = s.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11.3), Inches(1))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Quarterly results"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = DARK
rows, cols = 4, 3
tbl = s.shapes.add_table(rows, cols, Inches(1), Inches(2.3), Inches(10), Inches(3.5)).table
data = [["Quarter", "Revenue", "Growth"], ["Q1", "1.2M", "+8%"], ["Q2", "1.5M", "+25%"], ["Q3", "1.8M", "+20%"]]
for i, row in enumerate(data):
    for j, val in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = val
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(22)
        if i == 0:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

prs.save("public/sample-deck.pptx")
print("public/sample-deck.pptx written")
